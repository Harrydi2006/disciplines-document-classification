import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import configparser
from pathlib import Path
import threading
from typing import Dict, Optional
import json
import queue
import logging
import os
from concurrent.futures import ThreadPoolExecutor, as_completed
import time
import shutil

# 设置日志记录器
logger = logging.getLogger(__name__)

class QueueHandler(logging.Handler):
    def __init__(self, log_queue):
        super().__init__()
        self.log_queue = log_queue

    def emit(self, record):
        self.log_queue.put(self.format(record))

# 如果支持拖放，启用拖放功能
try:
    import windnd
    DRAG_DROP_SUPPORT = True
except ImportError:
    DRAG_DROP_SUPPORT = False
    logger.warning("未能导入 windnd 模块，拖放功能将被禁用")

class MainWindow:
    def __init__(self, config: configparser.ConfigParser):
        self.root = tk.Tk()
        self.root.title("文件分类助手")
        self.root.geometry("1200x900")
        
        self.config = config
        self.processing = False
        self.files_status: Dict[str, str] = {}
        self.files_results: Dict[str, str] = {}
        self.log_queue = queue.Queue()
        
        self._setup_logger()
        self._init_ui()
        self._center_window()
        self._start_log_monitor()
        
        # 如果支持拖放，启用拖放功能
        if DRAG_DROP_SUPPORT:
            windnd.hook_dropfiles(self.root, func=self._on_drop_files)
        
    def _setup_logger(self):
        """设置日志处理器"""
        try:
            # 创建日志目录
            log_dir = Path('logs')
            log_dir.mkdir(exist_ok=True)
            
            # 创建日志文件名（使用当前时间）
            timestamp = time.strftime("%Y%m%d_%H%M%S")
            log_file = log_dir / f'file_classifier_{timestamp}.log'
            
            # 尝试创建文件处理器，如果失败则使用备用文件名
            max_retries = 3
            retry_count = 0
            file_handler = None
            
            while retry_count < max_retries:
                try:
                    if retry_count > 0:
                        # 如果是重试，使用带有随机数的文件名
                        import random
                        random_suffix = random.randint(1000, 9999)
                        log_file = log_dir / f'file_classifier_{timestamp}_{random_suffix}.log'
                    
                    file_handler = logging.FileHandler(log_file, encoding='utf-8', mode='w', delay=True)
                    break
                except (PermissionError, OSError) as e:
                    retry_count += 1
                    if retry_count == max_retries:
                        # 如果所有重试都失败，使用标准错误输出
                        import sys
                        file_handler = logging.StreamHandler(sys.stderr)
                        logger.warning(f"无法创建日志文件，将使用标准错误输出: {str(e)}")
                    else:
                        time.sleep(0.1)  # 短暂延迟后重试
            
            file_handler.setLevel(logging.INFO)
            file_handler.setFormatter(
                logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s',
                                 datefmt='%Y-%m-%d %H:%M:%S')
            )
            
            # 创建队列处理器（用于GUI显示）
            self.queue_handler = QueueHandler(self.log_queue)
            self.queue_handler.setLevel(logging.INFO)
            self.queue_handler.setFormatter(
                logging.Formatter('%(asctime)s - %(levelname)s - %(message)s', 
                                 datefmt='%Y-%m-%d %H:%M:%S')
            )
            
            # 获取根日志记录器
            root_logger = logging.getLogger()
            root_logger.setLevel(logging.INFO)
            
            # 移除所有现有的处理器
            for handler in root_logger.handlers[:]:
                try:
                    handler.close()  # 尝试关闭现有的处理器
                except:
                    pass
                root_logger.removeHandler(handler)
            
            # 添加处理器
            root_logger.addHandler(file_handler)
            root_logger.addHandler(self.queue_handler)
            
            # 确保其他模块的日志也能显示
            logging.getLogger('file_classifier').propagate = True
            logging.getLogger('audio_processor').propagate = True
            
            # 记录启动信息
            logger.info("="*50)
            logger.info("程序启动")
            logger.info(f"日志文件: {log_file}")
            logger.info("="*50)
            
        except Exception as e:
            # 如果设置日志失败，至少尝试显示错误
            print(f"设置日志时出错: {str(e)}")
            messagebox.showerror("错误", f"设置日志时出错: {str(e)}")
        
    def _init_ui(self):
        # 创建主框架
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # 添加关于按钮到主窗口右上角
        about_button = ttk.Button(
            self.root,
            text="关于",
            command=self._show_about,
            width=8
        )
        about_button.pack(anchor=tk.NE, padx=10, pady=5)
        
        # 左侧：文件列表和控制按钮
        left_frame = ttk.Frame(main_frame)
        left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        # 路径设置
        path_frame = ttk.LabelFrame(left_frame, text="路径设置", padding="5")
        path_frame.pack(fill=tk.X, padx=5, pady=5)
        
        # 目标文件夹
        target_frame = ttk.Frame(path_frame)
        target_frame.pack(fill=tk.X, pady=2)
        ttk.Label(target_frame, text="目标文件夹:").pack(side=tk.LEFT)
        self.target_var = tk.StringVar(value=self.config['Paths']['target_folder'])
        target_entry = ttk.Entry(target_frame, textvariable=self.target_var)
        target_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=5)
        ttk.Button(
            target_frame,
            text="浏览",
            command=lambda: self._select_folder(self.target_var)
        ).pack(side=tk.LEFT)
        
        # 文件列表
        list_frame = ttk.LabelFrame(left_frame, text="文件列表", padding="5")
        list_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        
        # 添加按钮框架
        buttons_frame = ttk.Frame(list_frame)
        buttons_frame.pack(fill=tk.X, pady=(0, 5))
        
        # 文件操作按钮框架
        file_buttons_frame = ttk.Frame(buttons_frame)
        file_buttons_frame.pack(fill=tk.X, pady=(0, 5))
        
        # 添加文件和文件夹按钮
        ttk.Button(
            file_buttons_frame,
            text="添加文件",
            command=self._add_files
        ).pack(side=tk.LEFT, padx=2)
        
        ttk.Button(
            file_buttons_frame,
            text="添加文件夹",
            command=self._add_folder
        ).pack(side=tk.LEFT, padx=2)
        
        # 选择操作按钮框架
        select_buttons_frame = ttk.Frame(buttons_frame)
        select_buttons_frame.pack(fill=tk.X, pady=(0, 5))
        
        # 选择按钮
        ttk.Button(
            select_buttons_frame,
            text="全选",
            command=self._select_all
        ).pack(side=tk.LEFT, padx=2)
        
        ttk.Button(
            select_buttons_frame,
            text="反选",
            command=self._invert_selection
        ).pack(side=tk.LEFT, padx=2)
        
        ttk.Button(
            select_buttons_frame,
            text="取消选择",
            command=self._deselect_all
        ).pack(side=tk.LEFT, padx=2)
        
        # 文件类型选择按钮
        ttk.Button(
            select_buttons_frame,
            text="选择文档",
            command=lambda: self._select_by_type('doc')
        ).pack(side=tk.LEFT, padx=2)
        
        ttk.Button(
            select_buttons_frame,
            text="选择压缩包",
            command=lambda: self._select_by_type('archive')
        ).pack(side=tk.LEFT, padx=2)
        
        ttk.Button(
            select_buttons_frame,
            text="选择音频",
            command=lambda: self._select_by_type('audio')
        ).pack(side=tk.LEFT, padx=2)
        
        # 添加删除按钮
        ttk.Button(
            select_buttons_frame,
            text="删除选中",
            command=self._delete_selected
        ).pack(side=tk.LEFT, padx=2)
        
        # 创建树形视图
        columns = ("选择", "文件名", "状态", "分类结果")
        self.tree = ttk.Treeview(list_frame, columns=columns, show="headings")
        
        # 设置列
        self.tree.heading("选择", text="选择")
        self.tree.heading("文件名", text="文件名")
        self.tree.heading("状态", text="状态")
        self.tree.heading("分类结果", text="分类结果")
        
        self.tree.column("选择", width=50)
        self.tree.column("文件名", width=300)
        self.tree.column("状态", width=100)
        self.tree.column("分类结果", width=100)
        
        # 添加复选框状态字典
        self.checkboxes = {}
        
        # 绑定点击事件
        self.tree.bind('<Button-1>', self._on_click)
        
        # 添加滚动条
        scrollbar = ttk.Scrollbar(list_frame, orient=tk.VERTICAL, command=self.tree.yview)
        self.tree.configure(yscrollcommand=scrollbar.set)
        
        self.tree.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # 控制按钮
        control_frame = ttk.Frame(left_frame)
        control_frame.pack(fill=tk.X, pady=10)
        
        self.start_button = ttk.Button(
            control_frame,
            text="开始分类",
            command=self._start_classification,
            state='disabled'  # 初始状态为禁用
        )
        self.start_button.pack(side=tk.LEFT, padx=5)
        
        # 右侧：设置面板和日志
        right_frame = ttk.Frame(main_frame, width=400)
        right_frame.pack(side=tk.RIGHT, fill=tk.BOTH, padx=10)
        right_frame.pack_propagate(False)
        
        # 设置面板
        settings_frame = ttk.LabelFrame(right_frame, text="设置", padding="10")
        settings_frame.pack(fill=tk.X, expand=False)
        
        # 模型设置
        model_frame = ttk.LabelFrame(settings_frame, text="模型设置", padding="5")
        model_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(model_frame, text="模型名称:").pack(anchor=tk.W)
        self.model_name_var = tk.StringVar(value=self.config['Model']['model_name'])
        model_name_entry = ttk.Entry(model_frame, textvariable=self.model_name_var)
        model_name_entry.pack(fill=tk.X, pady=2)
        
        # API 设置
        api_frame = ttk.LabelFrame(settings_frame, text="API设置", padding="5")
        api_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(api_frame, text="API Key:").pack(anchor=tk.W)
        self.api_key_var = tk.StringVar(value=self.config['API']['api_key'])
        api_key_entry = ttk.Entry(api_frame, textvariable=self.api_key_var, show="*")
        api_key_entry.pack(fill=tk.X, pady=2)
        
        ttk.Label(api_frame, text="API Host:").pack(anchor=tk.W)
        self.api_host_var = tk.StringVar(value=self.config['API']['host'])
        api_host_entry = ttk.Entry(api_frame, textvariable=self.api_host_var)
        api_host_entry.pack(fill=tk.X, pady=2)
        
        # 添加测试按钮
        ttk.Button(
            api_frame,
            text="测试 API 连接",
            command=self._test_api
        ).pack(pady=5)
        
        # 提示词设置
        prompt_frame = ttk.LabelFrame(settings_frame, text="提示词设置", padding="5")
        prompt_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(prompt_frame, text="分类提示词:").pack(anchor=tk.W)
        self.prompt_var = tk.StringVar(value=self.config['Prompt']['classification_prompt'])
        prompt_entry = ttk.Entry(prompt_frame, textvariable=self.prompt_var)
        prompt_entry.pack(fill=tk.X, pady=2)
        
        # 功能设置
        features_frame = ttk.LabelFrame(settings_frame, text="功能设置", padding="5")
        features_frame.pack(fill=tk.X, pady=5)
        
        self.ocr_var = tk.BooleanVar(value=self.config.getboolean('Features', 'enable_ocr'))
        ttk.Checkbutton(features_frame, text="启用OCR", variable=self.ocr_var).pack(anchor=tk.W)
        
        self.audio_var = tk.BooleanVar(value=self.config.getboolean('Features', 'enable_audio'))
        ttk.Checkbutton(features_frame, text="启用音频识别", variable=self.audio_var).pack(anchor=tk.W)
        
        self.archive_var = tk.BooleanVar(value=self.config.getboolean('Features', 'enable_archive'))
        ttk.Checkbutton(features_frame, text="启用压缩包处理", variable=self.archive_var).pack(anchor=tk.W)
        
        # 添加子文件夹选项
        self.include_subfolders = tk.BooleanVar(value=self.config.getboolean('Features', 'include_subfolders', fallback=True))
        ttk.Checkbutton(features_frame, text="读取子文件夹", variable=self.include_subfolders).pack(anchor=tk.W)
        
        # 线程设置
        thread_frame = ttk.LabelFrame(settings_frame, text="线程设置", padding="5")
        thread_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(thread_frame, text="最大线程数:").pack(side=tk.LEFT)
        self.max_workers_var = tk.StringVar(value=self.config['Threading']['max_workers'])
        max_workers_entry = ttk.Entry(thread_frame, textvariable=self.max_workers_var, width=5)
        max_workers_entry.pack(side=tk.LEFT, padx=5)
        
        # 外部程序设置
        external_frame = ttk.LabelFrame(settings_frame, text="外部程序设置", padding="5")
        external_frame.pack(fill=tk.X, pady=5)
        
        # OCR 路径
        ocr_path_frame = ttk.Frame(external_frame)
        ocr_path_frame.pack(fill=tk.X, pady=2)
        ttk.Label(ocr_path_frame, text="Tesseract 路径:").pack(side=tk.LEFT)
        self.tesseract_path_var = tk.StringVar(value=self.config['OCR']['tesseract_path'])
        tesseract_entry = ttk.Entry(ocr_path_frame, textvariable=self.tesseract_path_var)
        tesseract_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=5)
        ttk.Button(
            ocr_path_frame,
            text="浏览",
            command=lambda: self._select_file(self.tesseract_path_var, "选择 Tesseract 可执行文件", "exe")
        ).pack(side=tk.LEFT)
        
        # FFmpeg 路径
        ffmpeg_path_frame = ttk.Frame(external_frame)
        ffmpeg_path_frame.pack(fill=tk.X, pady=2)
        ttk.Label(ffmpeg_path_frame, text="FFmpeg 路径:").pack(side=tk.LEFT)
        self.ffmpeg_path_var = tk.StringVar(value=self.config['Audio']['ffmpeg_path'])
        ffmpeg_entry = ttk.Entry(ffmpeg_path_frame, textvariable=self.ffmpeg_path_var)
        ffmpeg_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=5)
        ttk.Button(
            ffmpeg_path_frame,
            text="浏览",
            command=lambda: self._select_file(self.ffmpeg_path_var, "选择 FFmpeg 可执行文件", "exe")
        ).pack(side=tk.LEFT)
        
        # 保存按钮
        ttk.Button(
            settings_frame,
            text="保存设置",
            command=self._save_settings
        ).pack(pady=10)
        
        # 日志显示
        log_frame = ttk.LabelFrame(right_frame, text="日志", padding="5")
        log_frame.pack(fill=tk.BOTH, expand=True, pady=5)
        
        self.log_text = tk.Text(
            log_frame, 
            wrap=tk.WORD, 
            height=20  # 增加日志文本框高度
        )
        self.log_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        log_scrollbar = ttk.Scrollbar(log_frame, orient=tk.VERTICAL, command=self.log_text.yview)
        self.log_text.configure(yscrollcommand=log_scrollbar.set)
        log_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # 检查并更新按钮状态
        self._update_buttons_state()
        
        # 绑定 Ctrl+A 快捷键
        self.root.bind('<Control-a>', self._handle_ctrl_a)
        self.root.bind('<Control-A>', self._handle_ctrl_a)  # 大写的A
        
    def _select_folder(self, var: tk.StringVar):
        """选择文件夹"""
        folder = filedialog.askdirectory(
            initialdir=var.get() if var.get() else os.getcwd()
        )
        if folder:
            var.set(folder)
            # 保存设置
            self._save_settings()
            # 更新按钮状态
            self._update_buttons_state()
            
    def _select_file(self, var: tk.StringVar, title: str, file_type: str):
        filetypes = [
            ("所有支持的文件", "*.txt *.doc *.docx *.pdf *.ppt *.pptx *.jpg *.png *.jpeg *.zip *.rar *.7z *.mp3 *.wav *.m4a"),
            ("文档", "*.txt *.doc *.docx *.pdf *.ppt *.pptx"),  # 添加 ppt 和 pptx
            ("图片", "*.jpg *.png *.jpeg"),
            ("压缩包", "*.zip *.rar *.7z"),
            ("音频", "*.mp3 *.wav *.m4a"),
            ("所有文件", "*.*")
        ]
        filename = filedialog.askopenfilename(
            initialdir=str(Path(var.get()).parent),
            title=title,
            filetypes=filetypes
        )
        if filename:
            var.set(filename)
        
    def _start_log_monitor(self):
        """监控日志队列并更新显示"""
        def check_queue():
            while True:
                try:
                    record = self.log_queue.get_nowait()
                    self.log_text.insert(tk.END, record + '\n')
                    self.log_text.see(tk.END)
                    self.log_text.update_idletasks()  # 强制更新显示
                except queue.Empty:
                    break
            self.root.after(100, check_queue)
        
        check_queue()
        
    def _save_settings(self):
        """保存设置"""
        try:
            # 更新配置
            self.config['API']['api_key'] = self.api_key_var.get()
            self.config['API']['host'] = self.api_host_var.get()
            
            self.config['Model']['model_name'] = self.model_name_var.get()
            
            self.config['Paths']['target_folder'] = self.target_var.get()
            
            self.config['Prompt']['classification_prompt'] = self.prompt_var.get()
            
            self.config['Features']['enable_ocr'] = str(self.ocr_var.get())
            self.config['Features']['enable_audio'] = str(self.audio_var.get())
            self.config['Features']['enable_archive'] = str(self.archive_var.get())
            self.config['Features']['include_subfolders'] = str(self.include_subfolders.get())
            
            self.config['Threading']['max_workers'] = self.max_workers_var.get()
            
            self.config['OCR']['tesseract_path'] = self.tesseract_path_var.get()
            self.config['Audio']['ffmpeg_path'] = self.ffmpeg_path_var.get()
            
            # 保存到文件
            with open('config.conf', 'w', encoding='utf-8') as f:
                self.config.write(f)
                
            # 更新按钮状态
            self._update_buttons_state()
            
            messagebox.showinfo("成功", "设置已保存")
        except Exception as e:
            messagebox.showerror("错误", f"保存设置失败: {str(e)}")
            
    def _test_api(self):
        """测试 API 连接"""
        try:
            # 保存当前设置
            self._save_settings()
            
            # 创建临时分类器进行测试
            from file_classifier import FileClassifier
            classifier = FileClassifier()
            
            # 测试 API 调用
            test_content = "这是一个测试文本"
            result = classifier._call_api(test_content)
            
            messagebox.showinfo("测试成功", f"API 连接测试成功！\n返回结果: {result}")
        except Exception as e:
            messagebox.showerror("测试失败", f"API 连接测试失败：\n{str(e)}")

    def _check_config(self) -> bool:
        """检查配置是否完整"""
        # 检查 API 设置
        if not self.config['API']['api_key'] or not self.config['API']['host']:
            messagebox.showwarning("配置不完整", "请先配置 API 信息（主机地址和密钥）")
            return False
        
        if not self.config['Model']['model_name']:
            messagebox.showwarning("配置不完整", "请先配置模型名称")
            return False
        
        # 检查目标文件夹
        target_dir = Path(self.config['Paths']['target_folder'])
        if not target_dir.exists():
            messagebox.showwarning("配置不完整", "请先设置并创建目标文件夹")
            return False
        
        return True

    def _start_classification(self):
        """开始分类文件"""
        # 获取选中的文件
        selected_files = []
        for item in self.tree.get_children():
            if self.checkboxes.get(item, False):
                values = self.tree.item(item)['values']
                if values[2] != "已完成":  # 只处理未完成的文件
                    file_name = values[1]
                    # 从状态字典中获取完整路径
                    for file_str in self.files_status.keys():
                        if Path(file_str).name == file_name:
                            selected_files.append(file_str)
                            break
        
        if not selected_files:
            messagebox.showinfo("提示", "请先选择要处理的文件")
            return
        
        if self.processing:
            messagebox.showwarning("警告", "正在处理文件，请等待当前任务完成")
            return
        
        self.processing = True
        self.start_button.config(state='disabled')
        
        def process_thread():
            logger = logging.getLogger('file_classifier')
            
            try:
                # 重新加载配置
                self.config.read('config.conf', encoding='utf-8')
                
                # 检查目录
                target_dir = Path(self.config['Paths']['target_folder'])
                
                if not target_dir.exists():
                    logger.info(f"创建目标文件夹: {target_dir}")
                    target_dir.mkdir(parents=True, exist_ok=True)
                
                from file_classifier import FileClassifier
                classifier = FileClassifier()
                
                logger.info(f"开始处理 {len(selected_files)} 个文件")
                
                # 获取线程设置
                try:
                    max_workers = int(self.max_workers_var.get())
                    if max_workers <= 0:  # 自动分配模式
                        total_files = len(selected_files)
                        max_workers = min(12, (total_files + 29) // 30)  # 最多12个线程，每个线程处理30个文件
                    else:
                        max_workers = min(12, max_workers)  # 强制限制最大线程数为12
                except ValueError:
                    max_workers = 4  # 默认值
                
                with ThreadPoolExecutor(max_workers=max_workers) as executor:
                    # 创建任务列表，确保使用绝对路径
                    future_to_file = {
                        executor.submit(classifier.classify_file, Path(file_str).resolve()): file_str 
                        for file_str in selected_files
                    }
                    
                    # 处理完成的任务
                    for future in as_completed(future_to_file):
                        file_str = future_to_file[future]
                        file_name = Path(file_str).name
                        try:
                            # 更新状态为"处理中"
                            self.files_status[file_str] = "处理中"
                            self.root.after(0, self._update_file_status, file_name, "处理中", "")
                            
                            # 获取处理结果
                            subject, reason = future.result()
                            
                            # 更新状态和结果
                            self.files_status[file_str] = "已完成"
                            # 根据是否有原因来设置显示结果
                            if subject == "未知":
                                display_result = f"未知({reason})" if reason else "未知"
                            else:
                                display_result = subject
                            self.files_results[file_str] = display_result
                            self.root.after(0, self._update_file_status, file_name, "已完成", display_result)
                            
                            # 如果分类成功，移动文件
                            if subject != "未知":
                                target_path = target_dir / subject / file_name
                                target_path.parent.mkdir(parents=True, exist_ok=True)
                                
                                # 如果目标文件已存在，添加序号
                                counter = 1
                                while target_path.exists():
                                    new_name = f"{target_path.stem}_{counter}{target_path.suffix}"
                                    target_path = target_path.parent / new_name
                                    counter += 1
                                
                                # 移动文件
                                shutil.move(str(Path(file_str)), str(target_path))
                                logger.info(f"文件 {file_name} 已移动到 {subject} 目录")
                            
                        except Exception as e:
                            logger.error(f"处理文件 {file_name} 失败: {str(e)}")
                            self.files_status[file_str] = "失败"
                            self.files_results[file_str] = f"未知 (处理失败: {str(e)})"
                            self.root.after(0, self._update_file_status, file_name, "失败", self.files_results[file_str])
                
                logger.info("文件处理完成")
                # 最后再刷新一次文件列表
                self.root.after(0, self._refresh_tree_view)
                messagebox.showinfo("完成", "文件处理完成")
                
            except Exception as e:
                error_msg = f"处理文件时发生错误: {str(e)}"
                logger.error(error_msg)
                messagebox.showerror("错误", error_msg)
            finally:
                self.processing = False
                self.root.after(0, lambda: self.start_button.config(state='normal'))
        
        threading.Thread(target=process_thread, daemon=True).start()
        
    def _update_file_status(self, filename: str, status: str, result: str):
        """更新文件状态"""
        for item in self.tree.get_children():
            if self.tree.item(item)['values'][1] == filename:
                # 处理显示结果
                display_result = result
                if '(' in result:  # 如果结果中包含括号
                    subject, reason = result.split('(', 1)
                    subject = subject.strip()
                    if subject == '未知':
                        display_result = f"未知({reason}"  # 保持原有格式
                    else:
                        display_result = subject  # 只显示科目名
                
                self.tree.item(item, values=(self.tree.item(item)['values'][0], filename, status, display_result))
                break
        
        # 更新缓存
        for file_str in list(self.files_status.keys()):  # 使用 list 避免在迭代时修改
            if Path(file_str).name == filename:
                self.files_status[file_str] = status
                if result:  # 只在有结果时更新
                    self.files_results[file_str] = result
                
    def _center_window(self):
        """将窗口居中显示"""
        self.root.update_idletasks()
        width = self.root.winfo_width()
        height = self.root.winfo_height()
        x = (self.root.winfo_screenwidth() // 2) - (width // 2)
        y = (self.root.winfo_screenheight() // 2) - (height // 2)
        self.root.geometry(f'{width}x{height}+{x}+{y}')
        
    def _update_buttons_state(self):
        """更新按钮状态"""
        target_path = self.target_var.get()
        
        if target_path and Path(target_path).exists():
            self.start_button.config(state='normal')
        else:
            self.start_button.config(state='disabled')
            # 清空文件列表
            for item in self.tree.get_children():
                self.tree.delete(item)

    def _show_about(self):
        """显示关于对话框"""
        about_window = tk.Toplevel(self.root)
        about_window.title("关于文件分类助手")
        about_window.geometry("400x350")  # 增加关于窗口高度
        
        # 设置模态
        about_window.transient(self.root)
        about_window.grab_set()
        
        # 创建框架
        frame = ttk.Frame(about_window, padding="20")
        frame.pack(fill=tk.BOTH, expand=True)
        
        # 软件标题
        title_label = ttk.Label(
            frame,
            text="文件分类助手",
            font=("微软雅黑", 16, "bold")
        )
        title_label.pack(pady=(0, 10))
        
        # 版本信息
        version_label = ttk.Label(
            frame,
            text="版本 1.0.0",
            font=("微软雅黑", 10)
        )
        version_label.pack()
        
        # 软件描述
        description = (
            "文件分类助手是一个基于人工智能的文档分类工具，\n"
            "能够自动识别并分类不同学科的文档。\n\n"
            "支持的功能：\n"
            "• 文本文件分类\n"
            "• OCR 图片识别\n"
            "• 音频文件转写\n"
            "• 压缩包处理"
        )
        desc_label = ttk.Label(
            frame,
            text=description,
            justify=tk.CENTER,
            wraplength=350
        )
        desc_label.pack(pady=20)
        
        # 作者信息和 GitHub 链接
        author_frame = ttk.Frame(frame)
        author_frame.pack(pady=(10, 20))  # 增加上下边距
        
        # 加载 GitHub 图标
        try:
            # 创建并保存 GitHub 图标的 base64 字符串
            github_icon = """
            iVBORw0KGgoAAAANSUhEUgAAABAAAAAQCAYAAAAf8/9hAAAABHNCSVQICAgIfAhkiAAAAAlwSFlzAAAAdgAAAHYBTnsmCAAAABl0RVh0U29mdHdhcmUAd3d3Lmlua3NjYXBlLm9yZ5vuPBoAAADSSURBVDiNxZIxbsJAEEX/NyshCqcgBS2iQEEJBXfIBXKAXCGHyBFyAe5AQUGBUkSJkE2BIhcrcbFjr7WrxJGmGM3M/zPz96+Nqv4CeAZwB2DQtYEEPgG8A/gwAF4ADHE7pQC2AKbGOYcVgB2+x78H8Gic8/B7gdtpBbDQQ+P0zLG6vgWQmUMi4o1z+5/+1MnPyqW1XgGYA8gy9Oc3gIlzzn+rUNU+gFcAj1gvFxuPEuK99dMmJtLc2xBvrbUPYBeHlwBe1PJ0PLnFBs89gB9z7QlxfhBEIQAAAABJRU5ErkJggg==
            """
            github_image = tk.PhotoImage(data=github_icon)
            
            author_label = ttk.Label(
                author_frame,
                text="作者: "
            )
            author_label.pack(side=tk.LEFT)
            
            github_link = ttk.Label(
                author_frame,
                text="Harrydi2006",
                cursor="hand2",
                image=github_image,
                compound=tk.LEFT
            )
            github_link.image = github_image  # 保持引用
            github_link.pack(side=tk.LEFT)
            
            # 绑定点击事件
            def open_github(event):
                import webbrowser
                webbrowser.open("https://github.com/Harrydi2006")
            
            github_link.bind("<Button-1>", open_github)
            
        except Exception as e:
            # 如果图标加载失败，使用纯文本
            author_label = ttk.Label(
                author_frame,
                text="作者: Harrydi2006 (GitHub)",
                cursor="hand2"
            )
            author_label.pack()
            author_label.bind("<Button-1>", lambda e: webbrowser.open("https://github.com/Harrydi2006"))
        
        # 版权信息
        copyright_label = ttk.Label(
            frame,
            text="© 2024 All Rights Reserved",
            font=("微软雅黑", 8)
        )
        copyright_label.pack(side=tk.BOTTOM, pady=20)  # 增加底部边距
        
        # 居中窗口
        about_window.update_idletasks()
        width = about_window.winfo_width()
        height = about_window.winfo_height()
        x = (about_window.winfo_screenwidth() // 2) - (width // 2)
        y = (about_window.winfo_screenheight() // 2) - (height // 2)
        about_window.geometry(f'{width}x{height}+{x}+{y}')

    def _add_files(self):
        """添加文件到列表"""
        files = filedialog.askopenfilenames(
            title="选择文件",
            filetypes=[
                ("所有文件", "*.*"),  # 允许选择所有文件
                ("所有支持的文件", "*.txt *.doc *.docx *.pdf *.ppt *.pptx *.jpg *.png *.jpeg *.zip *.rar *.7z *.mp3 *.wav *.m4a"),
                ("文档", "*.txt *.doc *.docx *.pdf *.ppt *.pptx"),
                ("图片", "*.jpg *.png *.jpeg"),
                ("压缩包", "*.zip *.rar *.7z"),
                ("音频", "*.mp3 *.wav *.m4a")
            ]
        )
        
        for file in files:
            file_path = Path(file)
            if not self._is_file_in_tree(file_path):
                item = self.tree.insert('', tk.END, values=('☐', file_path.name, "等待处理", "未分类"))
                self.checkboxes[item] = False
                self.files_status[str(file_path)] = "等待处理"
                self.files_results[str(file_path)] = "未分类"

    def _add_folder(self):
        """添加文件夹中的文件到列表"""
        folder = filedialog.askdirectory(title="选择文件夹")
        if folder:
            folder_path = Path(folder)
            
            # 创建进度窗口
            progress_window = tk.Toplevel(self.root)
            progress_window.title("添加文件")
            progress_window.geometry("300x150")
            progress_window.transient(self.root)
            progress_window.grab_set()
            
            # 居中显示
            progress_window.update_idletasks()
            x = (progress_window.winfo_screenwidth() - progress_window.winfo_width()) // 2
            y = (progress_window.winfo_screenheight() - progress_window.winfo_height()) // 2
            progress_window.geometry(f"+{x}+{y}")
            
            # 创建进度条框架
            frame = ttk.Frame(progress_window, padding="20")
            frame.pack(fill=tk.BOTH, expand=True)
            
            # 状态标签
            status_label = ttk.Label(frame, text="正在扫描文件...", wraplength=250)
            status_label.pack(pady=(0, 10))
            
            # 进度条
            progress_var = tk.DoubleVar()
            progress_bar = ttk.Progressbar(
                frame,
                mode='determinate',
                variable=progress_var,
                maximum=100
            )
            progress_bar.pack(fill=tk.X, pady=(0, 10))
            
            # 计数标签
            count_label = ttk.Label(frame, text="")
            count_label.pack()
            
            # 创建更新队列
            update_queue = queue.Queue()
            
            def queue_update(**kwargs):
                """将更新添加到队列"""
                update_queue.put(kwargs)
            
            def update_ui_from_queue():
                """从队列中获取并处理UI更新"""
                if not progress_window.winfo_exists():
                    return
                try:
                    while True:
                        kwargs = update_queue.get_nowait()
                        if 'status' in kwargs:
                            status_label.config(text=kwargs['status'])
                        if 'progress' in kwargs:
                            progress_var.set(kwargs['progress'])
                        if 'count' in kwargs:
                            count_label.config(text=kwargs['count'])
                except queue.Empty:
                    pass
                finally:
                    if progress_window.winfo_exists():
                        self.root.after(100, update_ui_from_queue)
            
            def scan_folder():
                try:
                    # 第一阶段：统计文件数量
                    queue_update(
                        status="正在统计文件数量...",
                        count="请稍候..."
                    )
                    
                    total_files = 0
                    existing_count = 0
                    
                    if self.include_subfolders.get():
                        for root, _, files in os.walk(folder_path):
                            total_files += len(files)
                    else:
                        total_files = sum(1 for f in folder_path.iterdir() if f.is_file())
                    
                    # 第二阶段：处理文件
                    processed_count = 0
                    added_count = 0
                    
                    def process_file(file_path):
                        nonlocal processed_count, added_count, existing_count
                        
                        if not self._is_file_in_tree(file_path):
                            self.root.after(0, self._add_file_to_tree, file_path)
                            added_count += 1
                        else:
                            existing_count += 1
                        
                        processed_count += 1
                        progress = (processed_count / total_files) * 100
                        
                        if processed_count % 10 == 0 or processed_count == total_files:
                            queue_update(
                                status=f"正在添加文件... ({added_count} 个新文件)",
                                progress=progress,
                                count=f"已处理: {processed_count}/{total_files}\n"
                                      f"已添加: {added_count}\n"
                                      f"已存在: {existing_count}"
                            )
                    
                    if self.include_subfolders.get():
                        for root, _, files in os.walk(folder_path):
                            for file in files:
                                process_file(Path(root) / file)
                    else:
                        for file_path in folder_path.iterdir():
                            if file_path.is_file():
                                process_file(file_path)
                    
                    # 完成处理
                    queue_update(
                        status="处理完成",
                        progress=100,
                        count=f"总文件数: {total_files}\n"
                              f"新增文件: {added_count}\n"
                              f"已存在: {existing_count}"
                    )
                    
                    # 延迟关闭窗口
                    self.root.after(1500, progress_window.destroy)
                    
                    if added_count == 0:
                        self.root.after(0, lambda: messagebox.showinfo(
                            "提示",
                            f"未添加新文件\n"
                            f"总文件数: {total_files}\n"
                            f"已存在: {existing_count}"
                        ))
                    
                except Exception as e:
                    logger.error(f"扫描文件夹时出错: {str(e)}")
                    self.root.after(0, lambda: messagebox.showerror("错误", f"扫描文件夹时出错：\n{str(e)}"))
                    self.root.after(100, progress_window.destroy)
            
            # 启动UI更新循环
            update_ui_from_queue()
            
            # 在新线程中执行扫描
            threading.Thread(target=scan_folder, daemon=True).start()

    def _batch_add_files_to_tree(self, file_tuples):
        """批量添加文件到树形视图"""
        try:
            # 预先准备所有数据
            items_data = []
            for abs_path, name in file_tuples:
                # 检查文件是否已存在
                if any(Path(p).name == name for p in self.files_status.keys()):
                    continue
                
                # 添加到树形视图
                item = self.tree.insert('', tk.END, values=(
                    '☐', name, "等待处理", "未分类"
                ))
                
                # 更新状态字典
                self.checkboxes[item] = False
                self.files_status[abs_path] = "等待处理"
                self.files_results[abs_path] = "未分类"
                items_data.append((item, abs_path))
            
            # 如果有文件被添加，启用按钮
            if items_data:
                self.start_button.config(state='normal')
            
            return len(items_data)
            
        except Exception as e:
            logger.error(f"批量添加文件到树形视图时出错: {str(e)}")
            return 0

    def _process_file_batch(self, items, processed_files, total_files, queue_update, operation='add'):
        """处理文件批次"""
        batch_results = []
        last_update_time = time.time()
        
        if operation == 'add':
            # 批量处理添加操作
            file_tuples = []
            for item in items:
                try:
                    if not self._is_file_in_tree(item):
                        result = self._add_file_to_tree(item)
                        if result:
                            file_tuples.append(result)
                    
                    processed_files[0] += 1
                    progress = (processed_files[0] / total_files) * 100
                    
                    # 每100ms更新一次UI
                    current_time = time.time()
                    if current_time - last_update_time >= 0.1:
                        queue_update(
                            file=f"正在处理: {item.name}",
                            count=f"已处理: {processed_files[0]}/{total_files}",
                            progress=progress
                        )
                        last_update_time = current_time
                    
                except Exception as e:
                    logger.error(f"处理文件 {item} 时出错: {str(e)}")
                    continue
            
            # 批量添加到树形视图
            if file_tuples:
                self.root.after(0, self._batch_add_files_to_tree, file_tuples)
            
            return []  # 不再需要返回结果
        
        else:  # operation == 'delete'
            # 删除操作的代码保持不变
            ...

    def _add_file_to_tree(self, file_path):
        """添加文件到树形视图"""
        try:
            # 使用绝对路径存储
            abs_path = str(file_path.resolve())
            
            # 检查文件是否已存在
            if self._is_file_in_tree(file_path):
                return None
            
            # 添加到树形视图
            item = self.tree.insert('', tk.END, values=(
                '☐', file_path.name, "等待处理", "未分类"
            ))
            
            # 更新状态字典
            self.checkboxes[item] = False
            self.files_status[abs_path] = "等待处理"
            self.files_results[abs_path] = "未分类"
            
            # 启用开始按钮
            self.start_button.config(state='normal')
            
            return item
            
        except Exception as e:
            logger.error(f"添加文件到树形视图时出错: {str(e)}")
            return None

    def _is_file_in_tree(self, file_path: Path) -> bool:
        """检查文件是否已在列表中（使用绝对路径判断）"""
        try:
            file_abs_path = str(file_path.resolve())
            # 直接检查状态字典中是否存在该路径
            return file_abs_path in self.files_status
        except Exception as e:
            logger.warning(f"检查文件是否在树中时出错: {str(e)}")
            return False

    def _select_all(self):
        """全选"""
        for item in self.tree.get_children():
            if self.tree.item(item)['values'][2] != "已完成":
                self.checkboxes[item] = True
                self.tree.set(item, "选择", '☑')

    def _deselect_all(self):
        """取消全选"""
        for item in self.tree.get_children():
            self.checkboxes[item] = False
            self.tree.set(item, "选择", '☐')

    def _invert_selection(self):
        """反选"""
        for item in self.tree.get_children():
            if self.tree.item(item)['values'][2] != "已完成":
                self.checkboxes[item] = not self.checkboxes[item]
                self.tree.set(item, "选择", '☑' if self.checkboxes[item] else '☐')

    def _select_by_type(self, file_type: str):
        """根据文件类型选择"""
        extensions = {
            'doc': ('.txt', '.doc', '.docx', '.pdf', '.ppt', '.pptx'),
            'archive': ('.zip', '.rar', '.7z'),
            'audio': ('.mp3', '.wav', '.m4a')
        }
        
        for item in self.tree.get_children():
            values = self.tree.item(item)['values']
            if values[2] != "已完成":  # 只选择未完成的文件
                try:
                    file_name = str(values[1])  # 确保文件名是字符串
                    ext = Path(file_name).suffix.lower()
                    if ext in extensions[file_type]:
                        self.checkboxes[item] = True
                        self.tree.set(item, "选择", '☑')
                except Exception as e:
                    logger.error(f"处理文件类型选择时出错: {str(e)}")
                    continue

    def _on_click(self, event):
        """处理树形视图的点击事件"""
        region = self.tree.identify("region", event.x, event.y)
        if region == "cell":
            column = self.tree.identify_column(event.x)
            if column == "#1":  # 第一列（选择列）
                item = self.tree.identify_row(event.y)
                if item:  # 确保点击了有效的行
                    # 更新复选框状态
                    current_state = self.checkboxes.get(item, False)
                    self.checkboxes[item] = not current_state
                    # 更新显示
                    self.tree.set(item, "选择", "☑" if self.checkboxes[item] else "☐")

    def _handle_ctrl_a(self, event):
        """处理 Ctrl+A 快捷键"""
        self._select_all()
        return 'break'  # 阻止事件继续传播

    def _delete_selected(self):
        """删除选中的文件"""
        # 创建初始化进度窗口
        init_window = tk.Toplevel(self.root)
        init_window.title("准备删除")
        init_window.geometry("300x100")
        init_window.transient(self.root)
        init_window.grab_set()
        
        # 居中显示
        init_window.update_idletasks()
        x = (init_window.winfo_screenwidth() - init_window.winfo_width()) // 2
        y = (init_window.winfo_screenheight() - init_window.winfo_height()) // 2
        init_window.geometry(f"+{x}+{y}")
        
        # 添加提示标签
        init_label = ttk.Label(init_window, text="正在收集文件信息...", wraplength=250)
        init_label.pack(pady=20)
        
        # 创建事件和队列
        scan_complete = threading.Event()
        ui_update_queue = queue.Queue()
        deletion_ready = threading.Event()
        result_data = {'global_vars': None}  # 使用字典存储数据，避免作用域问题
        
        def update_scan_progress():
            """更新扫描进度"""
            logger.info("开始更新扫描进度")
            last_update_time = time.time()
            update_interval = 0.1  # 100ms更新一次
            
            try:
                while not scan_complete.is_set():
                    try:
                        msg = ui_update_queue.get(timeout=0.1)
                        current_time = time.time()
                        
                        if current_time - last_update_time >= update_interval:
                            if init_window.winfo_exists():
                                init_label.config(text=msg)
                                init_window.update_idletasks()
                                logger.debug(f"进度更新: {msg}")
                            last_update_time = current_time
                            
                    except queue.Empty:
                        if init_window.winfo_exists():
                            init_window.update_idletasks()
                        else:
                            logger.warning("进度窗口已关闭")
                            break
                
                logger.info("扫描进度更新完成")
                # 如果扫描完成且需要创建删除窗口
                if deletion_ready.is_set() and init_window.winfo_exists():
                    logger.info("准备创建删除窗口")
                    init_window.destroy()
                    # 使用更短的延迟，并传递数据
                    self.root.after(1, lambda: start_deletion(result_data['global_vars']))
                
            except Exception as e:
                error_msg = f"更新扫描进度时出错: {str(e)}"
                logger.error(error_msg, exc_info=True)
                if init_window.winfo_exists():
                    init_window.destroy()
                    self.root.after(0, lambda m=error_msg: messagebox.showerror("错误", m))
        
        def collect_items():
            """收集选中的项目"""
            try:
                logger.info("开始收集文件信息...")
                # 获取所有项目的ID列表
                all_items = self.tree.get_children()
                total_items = len(all_items)
                logger.info(f"列表中共有 {total_items} 个项目")
                
                if total_items == 0:
                    logger.info("列表为空，无需处理")
                    scan_complete.set()
                    self.root.after(0, lambda: messagebox.showinfo("提示", "列表中没有文件"))
                    return
                
                # 一次性获取所有选中状态
                selected_items = [item for item in all_items if self.checkboxes.get(item, False)]
                total_selected = len(selected_items)
                
                if not selected_items:
                    logger.info("未选择任何文件")
                    scan_complete.set()
                    self.root.after(0, lambda: messagebox.showinfo("提示", "请先选择要删除的文件"))
                    return
                
                logger.info(f"共选择了 {total_selected} 个文件")
                ui_update_queue.put(f"已选择 {total_selected} 个文件\n正在准备删除操作...")
                
                # 存储数据到共享字典中
                result_data['global_vars'] = {
                    'selected_items': selected_items,
                    'total': total_selected
                }
                
                # 标记准备就绪
                deletion_ready.set()
                scan_complete.set()
                logger.info("文件收集完成，准备开始删除操作")
                
            except Exception as e:
                error_msg = f"收集文件信息时出错: {str(e)}"
                logger.error(error_msg, exc_info=True)
                scan_complete.set()
                self.root.after(0, lambda m=error_msg: messagebox.showerror("错误", m))
        
        def start_deletion(vars_data):
            """创建并显示删除进度窗口"""
            if not vars_data:
                logger.error("没有要删除的数据")
                return
            
            try:
                progress_window = tk.Toplevel(self.root)
                progress_window.title("删除进度")
                progress_window.geometry("400x200")
                progress_window.transient(self.root)
                progress_window.grab_set()
                
                # 居中显示
                progress_window.update_idletasks()
                x = (progress_window.winfo_screenwidth() - progress_window.winfo_width()) // 2
                y = (progress_window.winfo_screenheight() - progress_window.winfo_height()) // 2
                progress_window.geometry(f"+{x}+{y}")
                
                # 进度条和标签
                progress_var = tk.DoubleVar()
                progress_bar = ttk.Progressbar(
                    progress_window,
                    variable=progress_var,
                    maximum=100
                )
                progress_bar.pack(fill=tk.X, padx=20, pady=10)
                
                status_label = ttk.Label(progress_window, text="准备删除...", wraplength=350)
                status_label.pack(pady=5)
                
                count_label = ttk.Label(progress_window, text="")
                count_label.pack(pady=5)
                
                # 创建删除操作的事件和队列
                delete_complete = threading.Event()
                delete_queue = queue.Queue()
                result_queue = queue.Queue()
                
                def process_batch(batch):
                    """处理一批文件"""
                    batch_results = []
                    batch_updates = set()
                    
                    try:
                        # 批量获取文件信息
                        batch_values = {item: self.tree.item(item)['values'] for item in batch}
                        
                        for item, values in batch_values.items():
                            if delete_complete.is_set():
                                return batch_results, batch_updates
                            
                            try:
                                if not values or len(values) < 2:
                                    continue
                                
                                file_name = str(values[1])
                                # 收集需要删除的文件路径，并确保使用Path对象
                                files_to_delete = set()
                                for file_str in self.files_status.keys():
                                    try:
                                        file_path = Path(file_str)
                                        if file_path.name == file_name:
                                            files_to_delete.add(file_str)
                                    except Exception as e:
                                        logger.error(f"处理文件路径时出错: {str(e)}")
                                        continue
                                
                                batch_updates.update(files_to_delete)
                                batch_results.append(item)
                                
                            except Exception as e:
                                logger.error(f"处理项目时出错: {str(e)}")
                        
                        return batch_results, batch_updates
                    
                    except Exception as e:
                        logger.error(f"批处理时出错: {str(e)}")
                        return [], set()
                
                def update_delete_progress():
                    """更新删除进度"""
                    try:
                        while not delete_complete.is_set():
                            try:
                                msg = delete_queue.get(timeout=0.1)
                                if isinstance(msg, dict):
                                    # 处理错误消息
                                    if 'error' in msg:
                                        self.root.after(0, lambda m=msg['error']: messagebox.showerror("错误", m))
                                        break
                                    # 处理正常进度更新
                                    if 'status' in msg:
                                        status_label.config(text=msg['status'])
                                    if 'progress' in msg:
                                        progress_var.set(msg['progress'])
                                    if 'count' in msg:
                                        count_label.config(text=msg['count'])
                                if progress_window.winfo_exists():
                                    progress_window.update_idletasks()
                            except queue.Empty:
                                if progress_window.winfo_exists():
                                    progress_window.update_idletasks()
                                else:
                                    break
                            
                            # 处理删除结果
                            try:
                                while True:
                                    items, updates = result_queue.get_nowait()
                                    if progress_window.winfo_exists():
                                        # 在主线程中执行UI更新
                                        self.root.after(0, lambda i=items, u=updates: self._batch_delete_items(i, u))
                            except queue.Empty:
                                pass
                        
                        # 删除完成后关闭窗口
                        if delete_complete.is_set() and progress_window.winfo_exists():
                            progress_window.destroy()
                            self.root.after(0, self._refresh_tree_view)
                    except Exception as e:
                        error_msg = f"更新删除进度时出错: {str(e)}"
                        logger.error(error_msg)
                        if progress_window.winfo_exists():
                            progress_window.destroy()
                            self.root.after(0, lambda m=error_msg: messagebox.showerror("错误", m))
                
                def delete_files():
                    """执行删除操作"""
                    try:
                        selected_items = vars_data['selected_items']
                        total = vars_data['total']
                        
                        # 根据文件数量动态调整批次大小和线程数
                        batch_size = min(max(total // 20, 100), 1000)  # 动态批次大小，最小100，最大1000
                        max_workers = min(max(total // 5000, 2), 8)    # 动态线程数，最小2，最大8
                        
                        logger.info(f"删除操作配置 - 批次大小: {batch_size}, 线程数: {max_workers}")
                        
                        # 创建批次
                        batches = [selected_items[i:i + batch_size] for i in range(0, total, batch_size)]
                        
                        processed = 0
                        last_update_time = time.time()
                        update_interval = 0.1  # 100ms更新一次
                        
                        delete_queue.put({
                            'status': "正在初始化删除操作...",
                            'progress': 0,
                            'count': f"已处理: 0/{total}"
                        })
                        
                        # 使用线程池处理
                        with ThreadPoolExecutor(max_workers=max_workers) as executor:
                            futures = [executor.submit(process_batch, batch) for batch in batches]
                            
                            for future in as_completed(futures):
                                if delete_complete.is_set():
                                    break
                                
                                try:
                                    items, updates = future.result()
                                    if items:
                                        result_queue.put((items, updates))
                                        processed += len(items)
                                        
                                        # 控制UI更新频率
                                        current_time = time.time()
                                        if current_time - last_update_time >= update_interval:
                                            progress = (processed / total) * 100
                                            delete_queue.put({
                                                'status': f"正在删除...\n已处理: {processed}/{total}",
                                                'progress': progress,
                                                'count': f"剩余: {total - processed}"
                                            })
                                            last_update_time = current_time
                                
                                except Exception as batch_error:
                                    error_msg = f"处理批次时出错: {str(batch_error)}"
                                    logger.error(error_msg)
                                    delete_queue.put({
                                        'status': f"错误: {error_msg}",
                                        'progress': (processed / total) * 100,
                                        'count': f"已处理: {processed}/{total}"
                                    })
                    
                    except Exception as e:
                        error_msg = f"删除文件时出错: {str(e)}"
                        logger.error(error_msg)
                        delete_queue.put({'error': error_msg})
                        delete_complete.set()
                
                # 启动删除进度更新线程
                threading.Thread(target=update_delete_progress, daemon=True).start()
                
                # 启动删除操作线程
                threading.Thread(target=delete_files, daemon=True).start()
                
                # 处理窗口关闭
                def on_closing():
                    delete_complete.set()
                    progress_window.destroy()
                
                progress_window.protocol("WM_DELETE_WINDOW", on_closing)
                
            except Exception as e:
                error_msg = f"创建删除窗口时出错: {str(e)}"
                logger.error(error_msg)
                messagebox.showerror("错误", error_msg)
        
        # 启动进度更新线程
        threading.Thread(target=update_scan_progress, daemon=True).start()
        
        # 启动文件收集线程
        collect_thread = threading.Thread(target=collect_items, daemon=True)
        collect_thread.start()

    def _batch_delete_items(self, items, updates):
        """在主线程中批量删除项目"""
        try:
            # 批量删除树形视图项目
            for item in items:
                if self.tree.exists(item):
                    self.tree.delete(item)
            
            # 批量更新状态字典
            for file_str in updates:
                self.files_status.pop(file_str, None)
                self.files_results.pop(file_str, None)
            
            # 批量更新复选框状态
            for item in items:
                self.checkboxes.pop(item, None)
            
            # 更新按钮状态
            remaining = len(self.tree.get_children())
            self.start_button.config(
                state='normal' if remaining > 0 else 'disabled'
            )
            
            # 强制更新UI
            self.tree.update_idletasks()
            self.root.update_idletasks()  # 确保主窗口也更新
            
        except Exception as e:
            logger.error(f"批量删除项目时出错: {str(e)}", exc_info=True)
            messagebox.showerror("错误", f"删除项目时出错：{str(e)}")

    def _refresh_tree_view(self):
        """刷新树形视图"""
        try:
            # 保存当前所有项目的信息
            items_info = []
            for item in self.tree.get_children():
                values = self.tree.item(item)['values']
                if values and len(values) >= 4:
                    items_info.append({
                        'values': values,
                        'checked': self.checkboxes.get(item, False)
                    })
            
            # 清空树形视图
            for item in self.tree.get_children():
                self.tree.delete(item)
            
            # 重新添加仍在状态字典中的项目
            for info in items_info:
                file_name = info['values'][1]
                # 只添加仍在状态字典中的文件
                if any(Path(p).name == file_name for p in self.files_status.keys()):
                    item = self.tree.insert('', tk.END, values=info['values'])
                    self.checkboxes[item] = info['checked']
            
            # 强制更新UI
            self.tree.update_idletasks()
            self.root.update_idletasks()
            
        except Exception as e:
            logger.error(f"刷新树形视图时出错: {str(e)}")

    def _on_drop_files(self, files):
        """处理文件拖放"""
        # windnd 返回的是字节字符串，需要解码
        files = [file.decode('gbk') for file in files]
        added_count = 0
        
        # 创建进度窗口
        progress_window = tk.Toplevel(self.root)
        progress_window.title("添加文件")
        progress_window.geometry("300x150")
        progress_window.transient(self.root)
        progress_window.grab_set()
        
        # 居中显示
        progress_window.update_idletasks()
        x = (progress_window.winfo_screenwidth() - progress_window.winfo_width()) // 2
        y = (progress_window.winfo_screenheight() - progress_window.winfo_height()) // 2
        progress_window.geometry(f"+{x}+{y}")
        
        # 创建进度条框架
        frame = ttk.Frame(progress_window, padding="20")
        frame.pack(fill=tk.BOTH, expand=True)
        
        # 状态标签
        status_label = ttk.Label(frame, text="正在添加文件...", wraplength=250)
        status_label.pack(pady=(0, 10))
        
        # 进度条
        progress_var = tk.DoubleVar()
        progress_bar = ttk.Progressbar(
            frame,
            mode='determinate',
            variable=progress_var,
            maximum=100
        )
        progress_bar.pack(fill=tk.X, pady=(0, 10))
        
        # 计数标签
        count_label = ttk.Label(frame, text="")
        count_label.pack()
        
        def process_files():
            nonlocal added_count
            try:
                total_files = len(files)
                processed_files = 0
                
                for file in files:
                    file_path = Path(file)
                    
                    # 如果是文件夹，递归处理
                    if file_path.is_dir():
                        if self.include_subfolders.get():
                            for sub_file in file_path.rglob('*'):
                                if sub_file.is_file() and not self._is_file_in_tree(sub_file):
                                    self.root.after(0, self._add_file_to_tree, sub_file)
                                    added_count += 1
                        else:
                            for sub_file in file_path.glob('*'):
                                if sub_file.is_file() and not self._is_file_in_tree(sub_file):
                                    self.root.after(0, self._add_file_to_tree, sub_file)
                                    added_count += 1
                    # 如果是文件，直接添加
                    elif file_path.is_file() and not self._is_file_in_tree(file_path):
                        self.root.after(0, self._add_file_to_tree, file_path)
                        added_count += 1
                    
                    processed_files += 1
                    progress = (processed_files / total_files) * 100
                    
                    # 更新进度
                    self.root.after(0, lambda p=progress, c=processed_files: (
                        progress_var.set(p),
                        count_label.config(text=f"已处理: {c}/{total_files}")
                    ))
                
                # 完成后显示结果
                self.root.after(100, lambda: messagebox.showinfo("完成", f"已添加 {added_count} 个文件"))
                self.root.after(200, progress_window.destroy)
                
            except Exception as e:
                self.root.after(0, lambda: messagebox.showerror("错误", f"添加文件时出错：\n{str(e)}"))
                self.root.after(100, progress_window.destroy)
        
        # 在新线程中处理文件
        threading.Thread(target=process_files, daemon=True).start()
        
        return "ok"  # 必须返回这个值以通知拖放系统操作完成

    def _classify_by_content(self, file_path: Path) -> tuple[str, Optional[str]]:
        """根据文件内容进行分类"""
        try:
            content = ""
            ext = file_path.suffix.lower()
            
            if ext == '.txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(2000)  # 只读取前2000个字符
                    
            elif ext in ['.doc', '.docx']:
                import docx
                doc = docx.Document(file_path)
                content = '\n'.join([p.text for p in doc.paragraphs])[:2000]
                
            elif ext == '.pdf':
                try:
                    from PyPDF2 import PdfReader
                    reader = PdfReader(file_path)
                    
                    # 检查是否加密
                    if reader.is_encrypted:
                        logger.warning(f"PDF文件已加密: {file_path}")
                        # 尝试使用OCR识别第一页
                        try:
                            import fitz  # PyMuPDF
                            import tempfile
                            from PIL import Image
                            
                            # 打开PDF文件
                            doc = fitz.open(file_path)
                            first_page = doc[0]
                            
                            # 将第一页转换为图片
                            pix = first_page.get_pixmap()
                            
                            # 创建临时文件保存图片
                            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                                pix.save(temp_img.name)
                                
                                # 使用OCR识别图片
                                subject, reason = self._classify_by_ocr(Path(temp_img.name))
                                
                            # 清理临时文件
                            try:
                                os.unlink(temp_img.name)
                            except:
                                pass
                                
                            return subject, reason
                            
                        except Exception as ocr_e:
                            logger.error(f"加密PDF的OCR识别失败: {str(ocr_e)}")
                            return '未知', "加密PDF的OCR识别失败"
                            
                    content = ''
                    # 尝试读取前两页
                    try:
                        for page in reader.pages[:2]:
                            try:
                                page_text = page.extract_text()
                                if page_text:
                                    content += page_text + '\n'
                            except Exception as e:
                                logger.error(f"PDF页面文本提取失败: {str(e)}")
                                continue
                                
                        content = content.strip()[:2000]  # 限制长度
                        
                        if not content:  # 如果没有提取到任何文本
                            logger.warning(f"未能从PDF提取到文本: {file_path}")
                            return '未知', "未能从PDF提取到文本"
                            
                    except Exception as e:
                        logger.error(f"PDF内容读取失败: {str(e)}")
                        return '未知', f"PDF内容读取失败: {str(e)}"
                        
                except Exception as e:
                    logger.error(f"PDF文件处理失败: {str(e)}")
                    return '未知', f"PDF文件处理失败: {str(e)}"
                
            elif ext in ['.ppt', '.pptx']:
                try:
                    if ext == '.ppt':  # 旧版本PPT
                        logger.warning(f"不支持的PPT格式(仅支持.pptx): {file_path}")
                        return '未知', "不支持的PPT格式(仅支持.pptx)"
                        
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    
                    texts = []
                    slide_count = 0
                    
                    # 只处理前5张幻灯片
                    for slide in prs.slides:
                        if slide_count >= 5:
                            break
                            
                        slide_texts = []
                        
                        try:
                            # 提取所有可能包含文本的元素
                            for shape in slide.shapes:
                                try:
                                    if hasattr(shape, "text") and shape.text.strip():
                                        slide_texts.append(shape.text.strip())
                                except Exception as e:
                                    logger.error(f"PPT形状文本提取失败: {str(e)}")
                                    continue
                                    
                            if slide_texts:
                                texts.extend(slide_texts)
                                slide_count += 1
                                
                        except Exception as e:
                            logger.error(f"PPT幻灯片处理失败: {str(e)}")
                            continue
                            
                    content = '\n'.join(texts)[:2000]
                    
                    if not content:  # 如果没有提取到任何文本
                        logger.warning(f"未能从PPT提取到文本: {file_path}")
                        return '未知', "未能从PPT提取到文本"
                        
                except Exception as e:
                    logger.error(f"PPT文件处理失败: {str(e)}")
                    return '未知', f"PPT文件处理失败: {str(e)}"
            
            if content.strip():
                return self._call_api(content), None
            return '未知', "文件内容为空"
            
        except Exception as e:
            logger.error(f"文件内容分类失败 {file_path}: {str(e)}", exc_info=True)
            return '未知', f"处理出错: {str(e)}"

    def run(self):
        self.root.mainloop() 