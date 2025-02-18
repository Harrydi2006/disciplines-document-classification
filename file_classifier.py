import configparser
import os
import logging
import json
import threading
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
import tqdm
import requests
from typing import Optional
import sys
from utils.logger import setup_logger
import time
import shutil
from tkinter import messagebox
import subprocess

# 添加必要的路径
if getattr(sys, 'frozen', False):
    # 运行于 PyInstaller 打包后的环境
    base_path = sys._MEIPASS
    if base_path not in sys.path:
        sys.path.insert(0, base_path)
else:
    # 运行于开发环境
    base_path = os.path.dirname(os.path.abspath(__file__))
    if base_path not in sys.path:
        sys.path.insert(0, base_path)

# 设置日志记录器
logger = setup_logger('file_classifier', 'file_classifier.log')

class APIError(Exception):
    """API调用错误的自定义异常"""
    pass

class FileClassifier:
    def __init__(self):
        logger.info("初始化文件分类器")
        
        # 线程池相关的属性
        self._executor = None
        self._allow_dynamic_threads = False
        self._current_content_threads = 0
        self._base_workers = 0
        self._max_additional_threads = 4
        self._threads = []
        
        # 先创建默认配置（如果需要）
        self._create_default_config()
        
        try:
            # 加载配置
            self.config = self._load_config()
            self.subjects = ['语文', '数学', '英语', '物理', '化学', '生物', '未知']
            
            # 在初始化时检查环境
            if not self._check_environment():
                raise EnvironmentError("环境检查失败，请重新运行安装程序")
            
            # 只创建目标文件夹结构
            self._setup_folders()
            
        except Exception as e:
            logger.error(f"初始化失败: {str(e)}")
            # 显示错误对话框
            messagebox.showerror("初始化失败", 
                "程序初始化失败，可能是因为：\n"
                "1. 配置文件不存在或损坏\n"
                "2. 必要的文件夹无法创建\n"
                "3. 外部程序未正确安装\n\n"
                "程序将重新进行首次运行设置。"
            )
            
            # 删除可能损坏的配置文件
            if Path('config.conf').exists():
                Path('config.conf').unlink()
            
            # 重新创建配置
            self._create_default_config()
            
            # 重新加载配置
            self.config = self._load_config()
            self.subjects = ['语文', '数学', '英语', '物理', '化学', '生物', '未知']
            self._setup_folders()

    def _create_default_config(self):
        """如果配置文件不存在，创建默认配置文件"""
        if not Path('config.conf').exists():
            config = configparser.ConfigParser()
            
            # 设置默认配置
            config['API'] = {
                'host': '',  # 留空，等待用户填写
                'api_key': ''  # 留空，等待用户填写
            }
            
            config['Model'] = {
                'model_name': ''  # 留空，等待用户填写
            }
            
            config['Paths'] = {
                'target_folder': ''   # 只保留目标文件夹
            }
            
            config['Prompt'] = {
                'classification_prompt': '请判断以下内容属于哪个学科（语文、数学、英语、物理、化学、生物）？如果无法判断，请回答"未知"。内容：'
            }
            
            config['Features'] = {
                'enable_ocr': 'true',
                'enable_audio': 'true',
                'enable_archive': 'true',
                'include_subfolders': 'true'  # 添加子文件夹选项
            }
            
            config['Threading'] = {
                'max_workers': '4'
            }
            
            config['OCR'] = {
                'tesseract_path': r'C:\Program Files\Tesseract-OCR\tesseract.exe'
            }
            
            config['Audio'] = {
                'ffmpeg_path': r'C:\ffmpeg\bin\ffmpeg.exe'
            }
            
            # 创建配置文件
            with open('config.conf', 'w', encoding='utf-8') as configfile:
                config.write(configfile)
            
            # 显示首次运行提示
            messagebox.showinfo("首次运行提示", 
                "欢迎使用文件分类助手！\n\n"
                "首次运行需要完成以下设置：\n"
                "1. 配置 API 信息（主机地址、密钥和模型）\n"
                "2. 设置源文件夹和目标文件夹\n"
                "3. 安装必要的外部程序\n\n"
                "点击确定开始配置。"
            )
            
            # 启动安装程序
            from gui.setup_window import SetupWindow
            setup = SetupWindow()
            setup.run()

    def _load_config(self) -> configparser.ConfigParser:
        config = configparser.ConfigParser()
        try:
            if not config.read('config.conf', encoding='utf-8'):
                raise FileNotFoundError("配置文件不存在")
            return config
        except Exception as e:
            logger.error(f"配置文件读取失败: {str(e)}", exc_info=True)
            raise

    def _setup_folders(self):
        """设置必要的文件夹结构"""
        try:
            # 只检查目标文件夹
            target_path = self.config['Paths']['target_folder']
            
            if not target_path:
                logger.warning("目标文件夹未设置")
                messagebox.showwarning(
                    "路径未设置",
                    "请先设置目标文件夹路径。\n"
                    "可以在主界面的路径设置中选择合适的目录。"
                )
                return False
            
            # 创建目标文件夹及其子目录
            base_dir = Path(target_path)
            
            # 创建目标文件夹及其子目录
            for subject in self.subjects:
                subject_dir = base_dir / subject
                subject_dir.mkdir(parents=True, exist_ok=True)
            
            logger.info("文件夹结构已设置完成")
            return True
            
        except Exception as e:
            logger.error(f"设置文件夹结构失败: {str(e)}")
            messagebox.showerror(
                "错误",
                f"创建文件夹结构失败：\n{str(e)}\n\n"
                "请确保有足够的权限创建文件夹。"
            )
            return False

    def classify_file(self, file_path: Path) -> tuple[str, Optional[str]]:
        try:
            # 首先通过文件名判断
            subject, reason = self._classify_by_filename(file_path)
            if subject != '未知':
                return subject, None

            # 获取文件扩展名
            ext = file_path.suffix.lower()
            
            # 定义支持的文件类型
            supported_types = {
                'document': ['.txt', '.doc', '.docx', '.pdf', '.ppt', '.pptx'],
                'image': ['.jpg', '.png', '.jpeg'],
                'archive': ['.zip', '.rar', '.7z'],
                'audio': ['.mp3', '.wav', '.m4a']
            }
            
            # 根据文件类型选择处理方法
            if ext in supported_types['document']:
                subject, reason = self._classify_by_content(file_path)
            elif (ext in supported_types['image'] and 
                  self.config.getboolean('Features', 'enable_ocr')):
                subject, reason = self._classify_by_ocr(file_path)
            elif (ext in supported_types['archive'] and 
                  self.config.getboolean('Features', 'enable_archive')):
                subject, reason = self._classify_archive(file_path)
            elif (ext in supported_types['audio'] and 
                  self.config.getboolean('Features', 'enable_audio')):
                subject, reason = self._classify_audio(file_path)
            else:
                # 对于不支持的文件类型，返回未知和原因
                reason = "不支持的文件类型或相关功能未启用"
                logger.info(f"{reason}: {file_path.name}")
                return '未知', reason
            
            return subject, reason
        except Exception as e:
            logger.error(f"文件分类失败 {file_path}: {str(e)}", exc_info=True)
            return '未知', f"处理出错: {str(e)}"

    def _call_api(self, content: str) -> str:
        try:
            headers = {
                'Authorization': f"Bearer {self.config['API']['api_key']}",
                'Content-Type': 'application/json'
            }
            
            data = {
                "model": self.config['Model']['model_name'],
                "messages": [
                    {"role": "system", "content": "你是一个文件分类助手。请从以下选项中选择一个：语文、数学、英语、物理、化学、生物、未知"},
                    {"role": "user", "content": f"{self.config['Prompt']['classification_prompt']}{content}"}
                ]
            }

            response = requests.post(
                f"{self.config['API']['host']}/chat/completions",
                headers=headers,
                json=data,
                timeout=30
            )
            
            if response.status_code != 200:
                error_msg = f"API调用失败: HTTP {response.status_code} - {response.text}"
                logger.error(error_msg)
                raise APIError(error_msg)

            try:
                response_data = response.json()
                logger.debug(f"API响应: {response_data}")
                
                if 'choices' not in response_data or not response_data['choices']:
                    raise APIError("API响应中没有choices字段")
                
                message_content = response_data['choices'][0]['message'].get('content', '')
                if not message_content:
                    raise APIError("API响应中没有content字段")
                
                # 从响应内容中提取学科
                for subject in self.subjects:
                    if subject in message_content:
                        return subject
                
                # 如果没有找到匹配的学科，返回未知
                return '未知'
                
            except json.JSONDecodeError as e:
                error_msg = f"API响应解析失败: {str(e)} - {response.text}"
                logger.error(error_msg)
                raise APIError(error_msg)
            except KeyError as e:
                error_msg = f"API响应格式错误: {str(e)} - {response.text}"
                logger.error(error_msg)
                raise APIError(error_msg)

        except requests.RequestException as e:
            error_msg = f"API请求异常: {str(e)}"
            logger.error(error_msg)
            raise APIError(error_msg)
        except Exception as e:
            error_msg = f"API调用过程中发生未知错误: {str(e)}"
            logger.error(error_msg)
            raise APIError(error_msg)

    def process_files(self):
        try:
            logger.info("开始处理文件")
            source_folder = Path(self.config['Paths']['source_folder'])
            if not source_folder.exists():
                logger.error(f"源文件夹不存在: {source_folder}")
                return

            files = list(source_folder.rglob('*'))
            files = [f for f in files if f.is_file()]
            
            if not files:
                logger.info("没有找到需要处理的文件")
                return

            # 创建分类缓存
            self._classification_cache = {}
            
            # 获取线程设置
            try:
                max_workers = int(self.config['Threading']['max_workers'])
                total_files = len(files)
                
                # 初始化线程池
                executor = self._initialize_thread_pool(max_workers, total_files)
                logger.info(f"开始处理 {total_files} 个文件，初始线程数: {self._base_workers}")
                
                try:
                    with tqdm.tqdm(total=total_files, desc="处理进度") as pbar:
                        futures = []
                        for file_path in files:
                            # 提交任务到线程池
                            future = executor.submit(self._process_single_file_with_thread_control, file_path)
                            futures.append(future)
                        
                        # 处理完成的任务
                        for future in as_completed(futures):
                            try:
                                file_path, subject = future.result()
                                if subject:
                                    # 更新缓存
                                    self._classification_cache[str(file_path)] = subject
                                    # 移动文件
                                    self._move_classified_file(file_path, subject)
                                pbar.update(1)
                            except Exception as e:
                                logger.error(f"处理文件失败: {str(e)}")
                                continue
                finally:
                    # 清理线程池
                    self._cleanup_thread_pool()

            except ValueError as e:
                logger.error(f"无效的线程数设置: {str(e)}")
                return
            except Exception as e:
                logger.error(f"处理文件时出错: {str(e)}")
                return

            logger.info("文件处理完成")

        except Exception as e:
            logger.error(f"处理文件过程中发生错误: {str(e)}", exc_info=True)
            raise

    def _process_single_file_with_thread_control(self, file_path: Path) -> tuple:
        """处理单个文件，包含线程控制逻辑"""
        try:
            # 检查是否需要读取文件内容
            ext = file_path.suffix.lower()
            content_processing = ext in ['.txt', '.doc', '.docx', '.pdf', '.ppt', '.pptx']
            
            if content_processing and self._allow_dynamic_threads:
                # 增加内容处理线程计数
                self._current_content_threads += 1
                
                # 如果所有线程都在处理内容，且未达到额外线程上限，添加新线程
                if (self._current_content_threads >= self._base_workers and 
                    len(self._threads) < self._base_workers + self._max_additional_threads):
                    self._add_thread()
            
            try:
                # 处理文件
                subject, reason = self.classify_file(file_path)
                return file_path, subject
            finally:
                if content_processing and self._allow_dynamic_threads:
                    # 减少内容处理线程计数
                    self._current_content_threads -= 1
                    
        except Exception as e:
            logger.error(f"处理文件失败 {file_path}: {str(e)}")
            return file_path, '未知'

    def _add_thread(self):
        """添加新的线程到线程池"""
        try:
            # 获取当前线程池
            executor = self._executor
            if executor and not executor._shutdown:
                # 增加最大线程数
                executor._max_workers += 1
                logger.info(f"增加线程池大小到 {executor._max_workers}")
        except Exception as e:
            logger.error(f"添加线程失败: {str(e)}")

    def _move_classified_file(self, file_path: Path, subject: str):
        """移动已分类的文件到目标目录"""
        try:
            # 确保subject不为空，如果为空则设为'未知'
            if not subject:
                subject = '未知'
                logger.info(f"文件 {file_path.name} 的分类为空，将移动到未知目录")
            
            target_dir = Path(self.config['Paths']['target_folder']) / subject
            target_path = target_dir / file_path.name
            
            # 确保目标目录存在
            target_dir.mkdir(parents=True, exist_ok=True)
            
            # 如果目标文件已存在，添加序号
            counter = 1
            while target_path.exists():
                new_name = f"{file_path.stem}_{counter}{file_path.suffix}"
                target_path = target_dir / new_name
                counter += 1
            
            # 移动文件
            shutil.move(str(file_path), str(target_path))
            logger.info(f"文件 {file_path.name} 已成功移动到 {subject} 目录")
            
        except Exception as e:
            logger.error(f"移动文件失败 {file_path}: {str(e)}")
            raise

    def _classify_by_filename(self, file_path: Path) -> tuple[str, Optional[str]]:
        try:
            filename = file_path.stem.lower()
            content = f"文件名：{filename}"
            subject = self._call_api(content)
            if subject != '未知':
                return subject, None

            # 获取文件扩展名
            ext = file_path.suffix.lower()
            
            # 定义支持的文件类型
            supported_types = {
                'document': ['.txt', '.doc', '.docx', '.pdf', '.ppt', '.pptx'],
                'image': ['.jpg', '.png', '.jpeg'],
                'archive': ['.zip', '.rar', '.7z'],
                'audio': ['.mp3', '.wav', '.m4a']
            }
            
            # 根据文件类型选择处理方法
            if ext in supported_types['document']:
                subject, reason = self._classify_by_content(file_path)
            elif (ext in supported_types['image'] and 
                  self.config.getboolean('Features', 'enable_ocr')):
                subject, reason = self._classify_by_ocr(file_path)
            elif (ext in supported_types['archive'] and 
                  self.config.getboolean('Features', 'enable_archive')):
                subject, reason = self._classify_archive(file_path)
            elif (ext in supported_types['audio'] and 
                  self.config.getboolean('Features', 'enable_audio')):
                subject, reason = self._classify_audio(file_path)
            else:
                # 对于不支持的文件类型，返回未知和原因
                reason = "不支持的文件类型或相关功能未启用"
                logger.info(f"{reason}: {file_path.name}")
                return '未知', reason
            
            return subject, reason
        except Exception as e:
            logger.error(f"文件名分类失败 {file_path}: {str(e)}", exc_info=True)
            return '未知', f"处理出错: {str(e)}"

    def _classify_by_content(self, file_path: Path) -> tuple[str, Optional[str]]:
        """根据文件内容进行分类"""
        try:
            content = ""
            ext = file_path.suffix.lower()
            
            if ext == '.txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(2000)  # 只读取前2000个字符
                    
            elif ext in ['.doc', '.docx']:
                import docx
                doc = docx.Document(file_path)
                content = '\n'.join([p.text for p in doc.paragraphs])[:2000]
                
            elif ext == '.pdf':
                try:
                    import fitz  # PyMuPDF
                    
                    # 尝试打开PDF文件
                    doc = fitz.open(file_path)
                    
                    # 检查是否加密
                    if doc.needs_pass:
                        logger.info(f"PDF文件已加密，尝试使用OCR: {file_path}")
                        try:
                            # 将第一页转换为图片
                            page = doc[0]
                            pix = page.get_pixmap()
                            img_path = file_path.with_suffix('.png')
                            pix.save(str(img_path))
                            
                            # 对图片进行OCR
                            result = self._classify_by_ocr(img_path)
                            
                            # 删除临时图片
                            if img_path.exists():
                                img_path.unlink()
                                
                            return result
                            
                        except Exception as e:
                            logger.error(f"加密PDF处理失败: {str(e)}")
                            return '未知', f"加密PDF处理失败: {str(e)}"
                    
                    # 如果PDF未加密，提取文本
                    content = ''
                    # 尝试读取前两页
                    for page_num in range(min(2, doc.page_count)):
                        try:
                            page = doc[page_num]
                            page_text = page.get_text()
                            if page_text:
                                content += page_text + '\n'
                        except Exception as e:
                            logger.error(f"PDF页面{page_num}文本提取失败: {str(e)}")
                            continue
                    
                    doc.close()
                    content = content.strip()[:2000]  # 限制长度
                    
                    if not content:  # 如果没有提取到任何文本
                        logger.warning(f"未能从PDF提取到文本，尝试OCR: {file_path}")
                        # 将第一页转换为图片并进行OCR
                        try:
                            doc = fitz.open(file_path)
                            page = doc[0]
                            pix = page.get_pixmap()
                            img_path = file_path.with_suffix('.png')
                            pix.save(str(img_path))
                            doc.close()
                            
                            result = self._classify_by_ocr(img_path)
                            
                            # 删除临时图片
                            if img_path.exists():
                                img_path.unlink()
                                
                            return result
                            
                        except Exception as e:
                            logger.error(f"PDF转图片失败: {str(e)}")
                            return '未知', f"PDF转图片失败: {str(e)}"
                    
                    return self._call_api(content), None
                    
                except Exception as e:
                    logger.error(f"PDF文件处理失败: {str(e)}")
                    return '未知', f"PDF文件处理失败: {str(e)}"
                
            elif ext in ['.ppt', '.pptx']:
                try:
                    if ext == '.ppt':  # 旧版本PPT
                        logger.warning(f"不支持的PPT格式(仅支持.pptx): {file_path}")
                        return '未知', "不支持的PPT格式(仅支持.pptx)"
                        
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    
                    texts = []
                    slide_count = 0
                    
                    # 只处理前5张幻灯片
                    for slide in prs.slides:
                        if slide_count >= 5:
                            break
                            
                        slide_texts = []
                        
                        try:
                            # 提取所有可能包含文本的元素
                            for shape in slide.shapes:
                                try:
                                    if hasattr(shape, "text") and shape.text.strip():
                                        slide_texts.append(shape.text.strip())
                                except Exception as e:
                                    logger.error(f"PPT形状文本提取失败: {str(e)}")
                                    continue
                                    
                            if slide_texts:
                                texts.extend(slide_texts)
                                slide_count += 1
                                
                        except Exception as e:
                            logger.error(f"PPT幻灯片处理失败: {str(e)}")
                            continue
                            
                    content = '\n'.join(texts)[:2000]  # 限制长度
                    
                    if not content:  # 如果没有提取到任何文本
                        logger.warning(f"未能从PPT提取到文本: {file_path}")
                        return '未知', "未能从PPT提取到文本"
                        
                except Exception as e:
                    logger.error(f"PPT文件处理失败: {str(e)}")
                    return '未知', f"PPT文件处理失败: {str(e)}"
            
            if content.strip():
                return self._call_api(content), None
            return '未知', "文件内容为空"
            
        except Exception as e:
            logger.error(f"文件内容分类失败 {file_path}: {str(e)}", exc_info=True)
            return '未知', f"处理出错: {str(e)}"

    def _classify_by_ocr(self, file_path: Path) -> tuple[str, Optional[str]]:
        try:
            import pytesseract
            from PIL import Image
            import subprocess
            
            logger.info(f"开始OCR识别: {file_path.name}")
            
            # 设置 Tesseract 路径
            tesseract_path = self.config['OCR']['tesseract_path']
            logger.info(f"使用 Tesseract 路径: {tesseract_path}")
            
            # 打开并预处理图像
            image = Image.open(file_path)
            
            # 使用 subprocess 直接调用 Tesseract
            try:
                result = subprocess.run(
                    [
                        str(tesseract_path),
                        'stdin',
                        'stdout',
                        '-l', 'chi_sim'
                    ],
                    input=image.tobytes(),
                    capture_output=True,
                    text=True,
                    creationflags=subprocess.CREATE_NO_WINDOW  # 隐藏窗口
                )
                
                if result.returncode == 0 and result.stdout:
                    text = result.stdout
                    # 清理OCR文本
                    text = ' '.join(text.split())
                    logger.info(f"OCR识别结果: {text[:200]}...")  # 只显示前200个字符
                    return self._call_api(text[:500]), None
                else:
                    logger.warning(f"OCR识别结果为空: {file_path.name}")
                    return '未知', "OCR识别结果为空"
                    
            except Exception as e:
                logger.error(f"OCR识别失败: {str(e)}")
                return '未知', f"OCR识别失败: {str(e)}"
            
            return '未知', "OCR识别未完成"
            
        except Exception as e:
            logger.error(f"OCR分类失败 {file_path}: {str(e)}", exc_info=True)
            return '未知', f"处理出错: {str(e)}"

    def _classify_archive(self, file_path: Path) -> tuple[str, Optional[str]]:
        try:
            import py7zr
            import rarfile
            import zipfile
            import tempfile
            
            # 创建临时目录
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)
                
                # 根据文件类型解压
                if file_path.suffix.lower() == '.zip':
                    with zipfile.ZipFile(file_path) as zip_ref:
                        file_list = zip_ref.namelist()
                        # 先尝试用文件名列表判断
                        content = "压缩包内文件：" + ", ".join(file_list)
                        result = self._call_api(content)
                        if result != '未知':
                            return result, None
                            
                        # 如果无法判断，解压第一个文本文件
                        for name in file_list:
                            if name.lower().endswith(('.txt', '.doc', '.docx', '.pdf')):
                                zip_ref.extract(name, temp_path)
                                return self._classify_by_content(temp_path / name)
                
                elif file_path.suffix.lower() == '.7z':
                    with py7zr.SevenZipFile(file_path, mode='r') as z:
                        file_list = z.getnames()
                        content = "压缩包内文件：" + ", ".join(file_list)
                        result = self._call_api(content)
                        if result != '未知':
                            return result, None
                            
                        for name in file_list:
                            if name.lower().endswith(('.txt', '.doc', '.docx', '.pdf')):
                                z.extract(temp_path, [name])
                                return self._classify_by_content(temp_path / name)
                
                elif file_path.suffix.lower() == '.rar':
                    with rarfile.RarFile(file_path) as rf:
                        file_list = rf.namelist()
                        content = "压缩包内文件：" + ", ".join(file_list)
                        result = self._call_api(content)
                        if result != '未知':
                            return result, None
                            
                        for name in file_list:
                            if name.lower().endswith(('.txt', '.doc', '.docx', '.pdf')):
                                rf.extract(name, temp_path)
                                return self._classify_by_content(temp_path / name)
            
            return '未知', "无法判断压缩包内容"
        except Exception as e:
            logger.error(f"压缩包分类失败 {file_path}: {str(e)}", exc_info=True)
            return '未知', f"处理出错: {str(e)}"

    def _classify_audio(self, file_path: Path) -> tuple[str, Optional[str]]:
        """分类音频文件"""
        try:
            logger.info(f"开始处理音频文件: {file_path.name}")
            # 首先尝试使用 vosk 进行离线识别
            try:
                from audio_processor import AudioProcessor
                processor = AudioProcessor()
                text = processor.transcribe_audio(file_path)
                if text:
                    logger.info(f"音频识别结果: {text[:200]}...")  # 只显示前200个字符
                    return self._call_api(text[:500]), None
            except ImportError:
                logger.warning("Vosk 模块未安装，将使用在线语音识别")
                
                # 使用 speech_recognition 作为备选
                import speech_recognition as sr
                from pydub import AudioSegment
                
                # 转换音频为 WAV 格式
                audio = AudioSegment.from_file(str(file_path))
                wav_path = file_path.with_suffix('.wav')
                audio.export(str(wav_path), format="wav")
                
                try:
                    # 使用 Google Speech Recognition
                    recognizer = sr.Recognizer()
                    with sr.AudioFile(str(wav_path)) as source:
                        audio = recognizer.record(source, duration=30)  # 只处理前30秒
                        text = recognizer.recognize_google(audio, language='zh-CN')
                        if text:
                            logger.info(f"在线语音识别结果: {text[:200]}...")  # 只显示前200个字符
                            return self._call_api(text[:500]), None
                finally:
                    # 清理临时文件
                    if wav_path.exists():
                        wav_path.unlink()
            
            logger.warning(f"音频识别结果为空: {file_path.name}")
            return '未知', "音频识别结果为空"
            
        except Exception as e:
            logger.error(f"音频分类失败 {file_path}: {str(e)}", exc_info=True)
            return '未知', f"处理出错: {str(e)}"

    def _check_environment(self) -> bool:
        """检查运行环境是否完整"""
        try:
            logger.info("开始检查运行环境...")
            
            # 检查 FFmpeg
            logger.info("检查 FFmpeg...")
            ffmpeg_path = Path(self.config['Audio']['ffmpeg_path'])
            if not ffmpeg_path.exists():
                logger.error("未找到 FFmpeg")
                return False
            
            # 测试 FFmpeg
            import subprocess
            try:
                result = subprocess.run(
                    [str(ffmpeg_path), '-version'],
                    capture_output=True,
                    text=True,
                    creationflags=subprocess.CREATE_NO_WINDOW  # 隐藏窗口
                )
                if result.returncode != 0:
                    logger.error("FFmpeg 测试失败")
                    return False
            except Exception as e:
                logger.error(f"FFmpeg 测试失败: {str(e)}")
                return False
            
            logger.info("FFmpeg 检查通过")
            
            # 检查 Tesseract
            logger.info("检查 Tesseract OCR...")
            tesseract_path = Path(self.config['OCR']['tesseract_path'])
            if not tesseract_path.exists():
                logger.error("未找到 Tesseract")
                return False
            
            # 测试 Tesseract
            try:
                result = subprocess.run(
                    [str(tesseract_path), '--version'],
                    capture_output=True,
                    text=True,
                    creationflags=subprocess.CREATE_NO_WINDOW  # 隐藏窗口
                )
                if result.returncode != 0:
                    logger.error("Tesseract 测试失败")
                    return False
                
                # 检查中文支持
                logger.info("检查 Tesseract 中文支持...")
                result = subprocess.run(
                    [str(tesseract_path), '--list-langs'],
                    capture_output=True,
                    text=True,
                    creationflags=subprocess.CREATE_NO_WINDOW  # 隐藏窗口
                )
                test_result = result.stdout
                
                if 'chi_sim' not in test_result:
                    logger.warning("未检测到中文支持，尝试安装...")
                    # 下载中文语言包的逻辑保持不变...
            
            except Exception as e:
                logger.error(f"Tesseract 测试失败: {str(e)}")
                return False
            
            logger.info("Tesseract 中文支持检查通过")
            
            # 检查配置文件
            logger.info("检查配置文件...")
            required_sections = ['API', 'Model', 'Paths', 'Threading', 'OCR', 'Audio']
            for section in required_sections:
                if section not in self.config:
                    logger.error(f"配置文件缺少 {section} 部分")
                    return False
            logger.info("配置文件检查通过")
            
            # 检查文件夹结构
            logger.info("检查文件夹结构...")
            target_dir = Path(self.config['Paths']['target_folder'])
            if not target_dir.exists():
                logger.warning("目标文件夹不存在")
                return False
            logger.info("文件夹结构检查通过")
            
            logger.info("环境检查完成：所有组件正常")
            return True
            
        except Exception as e:
            logger.error(f"环境检查失败: {str(e)}", exc_info=True)
            return False

    def _check_config(self) -> bool:
        """检查配置是否完整"""
        try:
            # 检查 API 设置
            if not self.config['API']['api_key']:
                return False
            
            # 只检查目标文件夹
            target_path = self.config['Paths']['target_folder']
            
            if not target_path:
                messagebox.showwarning(
                    "路径未设置",
                    "请先设置目标文件夹路径。\n"
                    "可以在主界面的路径设置中选择合适的目录。"
                )
                return False
            
            target_dir = Path(target_path)
            
            if not target_dir.exists():
                if not self._setup_folders():
                    return False
            
            # 检查外部程序
            if self.config.getboolean('Features', 'enable_ocr'):
                tesseract_path = Path(self.config['OCR']['tesseract_path'])
                if not tesseract_path.exists():
                    return False
            
            if self.config.getboolean('Features', 'enable_audio'):
                ffmpeg_path = Path(self.config['Audio']['ffmpeg_path'])
                if not ffmpeg_path.exists():
                    return False
            
            return True
            
        except Exception:
            return False

    def _initialize_thread_pool(self, max_workers: int, total_files: int):
        """初始化线程池"""
        if max_workers <= 0:  # 自动分配模式
            # 基础线程数：每30个文件1个线程，最少2个，最多12个
            self._base_workers = min(max(2, (total_files + 29) // 30), 12)
            max_workers = self._base_workers
            self._allow_dynamic_threads = True
            self._current_content_threads = 0
        else:
            max_workers = min(12, max_workers)  # 增加最大线程数限制到12
            self._base_workers = max_workers
            self._allow_dynamic_threads = False
        
        self._executor = ThreadPoolExecutor(max_workers=max_workers)
        self._threads = []
        return self._executor

    def _cleanup_thread_pool(self):
        """清理线程池"""
        if self._executor:
            self._executor.shutdown(wait=True)
            self._executor = None
            self._threads.clear()
            self._current_content_threads = 0

def main():
    try:
        # 设置日志
        log_dir = os.environ.get('LOG_DIR', 'logs')
        os.makedirs(log_dir, exist_ok=True)
        log_file = os.path.join(log_dir, 'file_classifier.log')
        
        # 配置日志记录
        logging.basicConfig(
            filename=log_file,
            level=logging.DEBUG,
            format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
        )
        
        logger = logging.getLogger('file_classifier')
        logger.info("程序启动")
        
        # 尝试关闭所有 tk 窗口
        if sys.platform == 'win32':
            try:
                import win32gui
                
                def close_tk_window(hwnd, extra):
                    classname = win32gui.GetClassName(hwnd)
                    if 'tk' in classname.lower():
                        win32gui.PostMessage(hwnd, 0x0010, 0, 0)  # WM_CLOSE
                
                win32gui.EnumWindows(close_tk_window, None)
            except ImportError:
                logger.warning("无法导入 win32gui 模块")
        
        # 创建 FileClassifier 实例
        classifier = FileClassifier()
        
        # 创建主窗口
        from gui.main_window import MainWindow
        window = MainWindow(classifier.config)
        
        # 运行主循环
        window.run()
        
    except Exception as e:
        print(f"Error: {str(e)}")
        if logger:
            logger.error(f"程序异常退出: {str(e)}", exc_info=True)
        sys.exit(1)

if __name__ == "__main__":
    main() 